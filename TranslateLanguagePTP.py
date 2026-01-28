import shutil
import os
import customtkinter as ctk
from tkinter import filedialog, messagebox
from pptx import Presentation
from openai import OpenAI

# =========================
# OpenAI Client
# =========================
client = OpenAI()

# =========================
# 翻訳関数
# =========================
def translate_text(text, source_lang, target_lang, model):
    prompt = f"""
次の文章を{source_lang}から{target_lang}に翻訳してください。
意味を変えず、自然な文章にしてください。

文章:
{text}
"""
    response = client.responses.create(
        model=model,
        input=prompt
    )
    return response.output_text.strip()

# =========================
# PPT 翻訳（完全保持版）
# =========================
def translate_ppt_inplace(
    input_path,
    output_path,
    source_lang,
    target_lang,
    model
):
    # 1. PPTコピー
    shutil.copy(input_path, output_path)

    prs = Presentation(output_path)

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            for paragraph in shape.text_frame.paragraphs:
                original_text = paragraph.text.strip()
                if not original_text:
                    continue

                translated = translate_text(
                    original_text,
                    source_lang,
                    target_lang,
                    model
                )

                # テキストのみ差し替え（レイアウト完全保持）
                paragraph.text = translated

    prs.save(output_path)

# =========================
# GUI
# =========================
class PPTTranslatorApp(ctk.CTk):

    def __init__(self): # アプリ起動時に一度だけ実行する ※ 基礎工事のようなもの
        super().__init__()
        self.output_dir = ""
        self.input_path = ""

        self.title("PPT Translator")
        self.geometry("520x420")
        ctk.set_appearance_mode("System")   # PCの設定に準ずる 他には…Light や Dark がある
        ctk.set_default_color_theme("blue") # アプリ全体のアクセントカラー ボタン や チェックボックス 選択時の色 等

        self.create_widgets()   # 画面の中身を実際に作成 ※ Label・Button・Entry・Frame の作成配置

    def create_widgets(self):
        TITLE_FONT = ("Meiryo", 20, "bold")
        SECTION_FONT = ("Segoe UI", 14, "bold")
        BODY_FONT = ("Segoe UI", 13)
        ctk.CTkLabel(self, text="PowerPoint 翻訳ツール", font=TITLE_FONT).pack(pady=15)

        # ▼ 1行分のレイアウト
        row = ctk.CTkFrame(self)
        row.pack(fill="x", padx=20, pady=8) # fill：横一杯に広げる 
        # 左
        ctk.CTkLabel( row, text="翻訳元PPT", width=90, anchor="w", font=BODY_FONT).pack(side="left", padx=(0, 10))
        # 中央
        self.input_label = ctk.CTkLabel( row, text="未選択", anchor="w", font=BODY_FONT)
        self.input_label.pack( side="left", fill="x", expand=True)
        # 右
        ctk.CTkButton( row, text="ファイル選択", width=90, command=self.select_input, font=BODY_FONT).pack(side="right")


       # ▼ 1行分のレイアウト
        row = ctk.CTkFrame(self)
        row.pack(fill="x", padx=20, pady=8)
        # 左
        ctk.CTkLabel( row, text="保存先", width=90, anchor="w", font=BODY_FONT).pack(side="left", padx=(10, 10))
        # 中央
        self.output_label = ctk.CTkLabel( row, text="未選択", anchor="w", font=BODY_FONT, width=200)
        self.output_label.pack( side="left", expand=True)
        # 右
        ctk.CTkButton( row, text="フォルダ選択", width=90, command=self.select_output_dir, font=BODY_FONT).pack(side="right")

        row = ctk.CTkFrame(self)
        row.pack(padx=20, pady=8)

        ctk.CTkLabel( row, text="翻訳前言語").pack(side="left", padx=10)

        self.source_lang = ctk.CTkEntry(row)
        self.source_lang.pack(pady=8, side="right", padx=(0, 10))
        self.source_lang.insert(0, "日本語")

        row = ctk.CTkFrame(self)
        row.pack(padx=20, pady=8)

        ctk.CTkLabel( row, text="翻訳後言語").pack(side="left", padx=10)

        self.target_lang = ctk.CTkEntry(row)
        self.target_lang.pack(pady=8, side="right", padx=(0, 10))
        self.target_lang.insert(0, "英語")
        # self.target_lang = ctk.CTkEntry(self)
        # self.target_lang.pack(pady=8)
        # self.target_lang.insert(0, "英語")

        self.model_entry = ctk.CTkEntry(self)
        self.model_entry.pack(pady=8)
        self.model_entry.insert(0, "gpt-4.1-mini")

        ctk.CTkButton(
            self,
            text="翻訳実行",
            command=self.run_translation,
            fg_color="green"
        ).pack(pady=20)

    def select_input(self):
        path = filedialog.askopenfilename(filetypes=[("PowerPoint", "*.pptx")])
        if path:
            self.input_path = path
            self.input_label.configure(text=os.path.basename(path))

    def select_output_dir(self):
        path = filedialog.askdirectory()
        if path:
            self.output_dir = path
            self.output_label.configure(text=path)

    def run_translation(self):
        if not self.input_path or not self.output_dir:
            messagebox.showerror("エラー", "入力PPTと保存先フォルダを選択してください")
            return

        try:
            base_name = os.path.splitext(os.path.basename(self.input_path))[0]
            output_path = os.path.join(
                self.output_dir,
                f"{base_name}_translated.pptx"
            )

            translate_ppt_inplace(
                self.input_path,
                output_path,
                self.source_lang.get(),
                self.target_lang.get(),
                self.model_entry.get()
            )

            messagebox.showinfo(
                "完了",
                f"翻訳が完了しました！\n保存先:\n{output_path}"
            )

        except Exception as e:
            messagebox.showerror("エラー", str(e))

if __name__ == "__main__":
    app = PPTTranslatorApp()
    app.mainloop()
